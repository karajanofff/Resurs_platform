import math
import re
from collections import Counter
from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


STOPWORDS = {
    "va",
    "yoki",
    "uchun",
    "bilan",
    "ham",
    "bu",
    "shu",
    "ning",
    "lar",
    "dan",
    "ga",
    "da",
    "the",
    "and",
    "for",
    "with",
}


def extract_text_from_file(file_path: str | Path) -> str:
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return "\n".join(page.extract_text() or "" for page in PdfReader(path).pages)
    if suffix == ".docx":
        return "\n".join(p.text for p in Document(path).paragraphs)
    if suffix == ".pptx":
        presentation = Presentation(path)
        return "\n".join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
    if suffix == ".txt":
        return path.read_text(encoding="utf-8", errors="ignore")
    raise ValueError("Faqat PDF, DOCX, PPTX yoki TXT fayllar qo'llab-quvvatlanadi.")


def clean_text(text: str) -> str:
    lowered = text.lower()
    lowered = re.sub(r"[^a-zA-Z0-9'`\s\u0400-\u04ff]", " ", lowered)
    return re.sub(r"\s+", " ", lowered).strip()


def extract_keywords(text: str, limit: int = 8) -> list[str]:
    tokens = [token for token in clean_text(text).split() if len(token) > 3 and token not in STOPWORDS]
    return [word for word, _ in Counter(tokens).most_common(limit)]


def split_text_into_sections(text: str, limit: int = 6) -> list[str]:
    paragraphs = [re.sub(r"\s+", " ", part).strip() for part in re.split(r"\n\s*\n+", text) if part.strip()]
    if len(paragraphs) >= 2:
        return paragraphs[:limit]

    sentences = [part.strip() for part in re.split(r"(?<=[.!?])\s+", text) if part.strip()]
    if len(sentences) <= 2:
        return sentences or [text.strip()]

    chunk_size = max(1, math.ceil(len(sentences) / min(limit, len(sentences))))
    return [" ".join(sentences[index : index + chunk_size]) for index in range(0, len(sentences), chunk_size)][:limit]


def calculate_similarity(resource_text: str, topic_text: str) -> float:
    documents = [clean_text(resource_text).split(), clean_text(topic_text).split()]
    vocabulary = sorted(set(documents[0]) | set(documents[1]))
    if not vocabulary:
        return 0.0

    vectors: list[list[float]] = []
    total_documents = len(documents)
    for document in documents:
        counts = Counter(document)
        length = max(len(document), 1)
        vector = []
        for term in vocabulary:
            tf = counts[term] / length
            docs_with_term = sum(1 for candidate in documents if term in candidate)
            idf = math.log((1 + total_documents) / (1 + docs_with_term)) + 1
            vector.append(tf * idf)
        vectors.append(vector)

    dot_product = sum(a * b for a, b in zip(vectors[0], vectors[1]))
    magnitude_a = math.sqrt(sum(value * value for value in vectors[0]))
    magnitude_b = math.sqrt(sum(value * value for value in vectors[1]))
    if not magnitude_a or not magnitude_b:
        return 0.0
    return round((dot_product / (magnitude_a * magnitude_b)) * 100, 2)


def generate_recommendation(score: float, keywords: list[str]) -> tuple[str, str]:
    if score >= 75:
        return "Mos", "Ushbu resurs aniqlangan fanga yuqori darajada mos keladi."
    if score >= 40:
        return "Qisman mos", "Resurs aniqlangan fanga qisman mos, lekin ayrim bo'limlari boshqa fanlarga yaqin."
    return "Mos emas", "Resursning umumiy mazmuni bitta fanga yetarli darajada yaqin emas."
