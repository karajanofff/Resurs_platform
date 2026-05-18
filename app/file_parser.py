from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


def extract_text_from_pdf(file_path: str | Path) -> str:
    return "\n".join(page.extract_text() or "" for page in PdfReader(file_path).pages)


def extract_text_from_docx(file_path: str | Path) -> str:
    return "\n".join(paragraph.text for paragraph in Document(file_path).paragraphs)


def extract_text_from_pptx(file_path: str | Path) -> str:
    presentation = Presentation(file_path)
    return "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


def extract_text_from_txt(file_path: str | Path) -> str:
    return Path(file_path).read_text(encoding="utf-8", errors="ignore")


def extract_text_from_file(file_path: str | Path) -> str:
    suffix = Path(file_path).suffix.lower()
    if suffix == ".pdf":
        return extract_text_from_pdf(file_path)
    if suffix == ".docx":
        return extract_text_from_docx(file_path)
    if suffix == ".pptx":
        return extract_text_from_pptx(file_path)
    if suffix == ".txt":
        return extract_text_from_txt(file_path)
    raise ValueError("Faqat PDF, DOCX, PPTX yoki TXT fayllar qo'llab-quvvatlanadi.")
